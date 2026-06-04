# pptx_translate.py
# PPTX translation engine — shape/XML logic ported from DocuTranslate
# (github.com/QinHan/DocuTranslate, MPL-2.0), translation calls replaced with
# direct Ollama integration.
#
# Public API:
#   translate_pptx(content, lang_to, translate_batch_fn, stop_event, chunk_size,
#                  insert_mode, separator)
#     -> bytes  (translated .pptx)

import io
import logging
import threading
from typing import Callable, Dict, List, Tuple

import regex
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.text.text import _Paragraph, TextFrame

log = logging.getLogger("pptx_translate")

# ---------------------------------------------------------------------------
# Language tag helper
# ---------------------------------------------------------------------------

_CJK_PATTERN = regex.compile(r'[\p{Han}\p{Hiragana}\p{Katakana}\p{Hangul}]')

_LANG_MAP = {
    "ukrainian": "uk-UA", "english": "en-US", "german": "de-DE",
    "french": "fr-FR", "spanish": "es-ES", "russian": "ru-RU",
    "polish": "pl-PL", "czech": "cs-CZ", "slovak": "sk-SK",
    "chinese": "zh-CN", "japanese": "ja-JP", "korean": "ko-KR",
    "arabic": "ar-SA", "hebrew": "he-IL", "turkish": "tr-TR",
    "italian": "it-IT", "portuguese": "pt-PT", "dutch": "nl-NL",
}


def _guess_lang_tag(lang_name: str, text: str) -> str:
    clean = lang_name.lower().strip()
    if clean in _LANG_MAP:
        return _LANG_MAP[clean]
    if regex.match(r'^[a-z]{2,3}(-[a-z0-9]+)?$', clean):
        return lang_name
    if _CJK_PATTERN.search(text):
        return "zh-CN"
    return "en-US"


# ---------------------------------------------------------------------------
# Style signature helpers
# ---------------------------------------------------------------------------

def _get_style_sig(run) -> tuple:
    rPr = run._r.rPr
    if rPr is None:
        return ("DEFAULT",)

    def bool_attr(tag_name):
        node = rPr.find(qn(f'a:{tag_name}'))
        if node is None:
            return None
        val = node.get('val')
        return val if val is not None else '1'

    u_node = rPr.find(qn('a:u'))
    latin = rPr.find(qn('a:latin'))
    color_sig = "INHERITED"
    for tag in ('solidFill', 'gradFill', 'noFill'):
        if rPr.find(qn(f'a:{tag}')) is not None:
            color_sig = tag
            break

    return (
        bool_attr('b'), bool_attr('i'),
        u_node.get('val') if u_node is not None else None,
        rPr.get('sz'),
        latin.get('typeface') if latin is not None else None,
        color_sig,
    )


def _same_style(run1, run2) -> bool:
    if _get_style_sig(run1) != _get_style_sig(run2):
        return False
    try:
        r1, r2 = run1._r, run2._r
        parent = r1.getparent()
        if parent != r2.getparent():
            return False
        if parent.index(r2) != parent.index(r1) + 1:
            return False
    except Exception:
        return False
    return True


# ---------------------------------------------------------------------------
# Segment collection
# ---------------------------------------------------------------------------

def _process_paragraph(paragraph: _Paragraph, elements: List[Dict], texts: List[str]):
    if not paragraph.runs:
        return
    state: Dict = {'current_runs': []}

    def flush():
        runs = state['current_runs']
        if not runs:
            return
        full_text = "".join(r.text for r in runs)
        if full_text.strip():
            elements.append({
                "runs": list(runs),
                "paragraph": paragraph,
                "text_frame": paragraph._parent,
            })
            texts.append(full_text)
        runs.clear()

    for run in paragraph.runs:
        if not run.text:
            continue
        last = state['current_runs'][-1] if state['current_runs'] else None
        if last and not _same_style(last, run):
            flush()
        state['current_runs'].append(run)
    flush()


def _process_text_frame(text_frame: TextFrame, elements: List[Dict], texts: List[str]):
    for paragraph in text_frame.paragraphs:
        _process_paragraph(paragraph, elements, texts)


def _process_shape(shape, elements: List[Dict], texts: List[str]):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _process_shape(child, elements, texts)
        return
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame") and cell.text_frame:
                    _process_text_frame(cell.text_frame, elements, texts)
        return
    if shape.has_text_frame:
        try:
            _process_text_frame(shape.text_frame, elements, texts)
        except Exception:
            pass


def _scan_deep_xml(slide_element, elements: List[Dict], texts: List[str]):
    """Catch text inside AlternateContent blocks (e.g. SmartArt fallbacks)."""
    MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    MC_ALT = f"{{{MC_NS}}}AlternateContent"
    MC_CHOICE = f"{{{MC_NS}}}Choice"

    for alt in slide_element.iter(MC_ALT):
        choice = alt.find(MC_CHOICE)
        if choice is None:
            continue
        for sp in choice.iter(qn('p:sp')):
            txBody = sp.find(qn('p:txBody'))
            if txBody is not None:
                try:
                    tf = TextFrame(txBody, None)
                    _process_text_frame(tf, elements, texts)
                except Exception as e:
                    log.warning(f"Deep XML scan error: {e}")


def _collect_all_segments(prs: Presentation) -> Tuple[List[Dict], List[str]]:
    elements: List[Dict] = []
    texts: List[str] = []

    def scan_slide_obj(slide_obj):
        for shape in slide_obj.shapes:
            _process_shape(shape, elements, texts)
        _scan_deep_xml(slide_obj.element, elements, texts)

    for slide in prs.slides:
        scan_slide_obj(slide)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            _process_text_frame(slide.notes_slide.notes_text_frame, elements, texts)

    for master in prs.slide_masters:
        scan_slide_obj(master)
        for layout in master.slide_layouts:
            scan_slide_obj(layout)

    return elements, texts


# ---------------------------------------------------------------------------
# Apply translations
# ---------------------------------------------------------------------------

def _apply_translation(info: Dict, original: str, translated: str,
                       lang_to: str, insert_mode: str, separator: str):
    runs = info["runs"]
    if not runs:
        return

    if insert_mode == "append":
        final_text = original + separator + translated
    elif insert_mode == "prepend":
        final_text = translated + separator + original
    else:
        final_text = translated

    primary = runs[0]
    try:
        primary.text = final_text
        rPr = primary._r.get_or_add_rPr()
        lang_tag = _guess_lang_tag(lang_to, final_text)
        rPr.set('lang', lang_tag)
        rPr.set('altLang', lang_tag)

        text_frame = info.get("text_frame")
        if text_frame and hasattr(text_frame, 'auto_size'):
            if text_frame.auto_size == MSO_AUTO_SIZE.NONE:
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception as e:
        log.warning(f"Error applying PPTX translation: {e}")
        return

    for run in runs[1:]:
        run.text = ""


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def translate_pptx(content: bytes,
                   lang_to: str,
                   translate_batch_fn: Callable,
                   stop_event: threading.Event,
                   chunk_size: int = 3000,
                   insert_mode: str = "replace",
                   separator: str = "\n") -> bytes:
    """Translate a PPTX file, preserving formatting.

    Args:
        content: raw .pptx bytes
        lang_to: target language name, e.g. "Ukrainian"
        translate_batch_fn: fn(batch: dict[str,str], lang_to: str, stop_event) -> dict[str,str]
        stop_event: set to abort mid-translation
        chunk_size: max characters per LLM request
        insert_mode: "replace" | "append" | "prepend"
        separator: string between original and translation in append/prepend mode

    Returns:
        Translated .pptx bytes.
    """
    prs = Presentation(io.BytesIO(content))
    elements, originals = _collect_all_segments(prs)

    if not originals:
        log.info("No translatable text found in PPTX")
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    translated: List[str] = [""] * len(originals)
    batch_ids: Dict[str, int] = {}
    current_batch: Dict[str, str] = {}
    current_chars = 0

    def flush_batch():
        if not current_batch or stop_event.is_set():
            return
        result = translate_batch_fn(current_batch, lang_to, stop_event)
        for k, v in result.items():
            idx = batch_ids.get(k)
            if idx is not None:
                translated[idx] = v or originals[idx]

    for i, text in enumerate(originals):
        if stop_event.is_set():
            break
        key = str(i)
        if current_chars + len(text) > chunk_size and current_batch:
            flush_batch()
            batch_ids.clear()
            current_batch.clear()
            current_chars = 0
        current_batch[key] = text
        batch_ids[key] = i
        current_chars += len(text)

    flush_batch()

    for i, t in enumerate(translated):
        if not t:
            translated[i] = originals[i]

    for info, orig, trans in zip(elements, originals, translated):
        _apply_translation(info, orig, trans, lang_to, insert_mode, separator)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
